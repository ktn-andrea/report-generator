#!/usr/bin/env python3.9

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from typing import List
import numpy as np
import logging
from os import path
from pathlib import Path
logger = logging.getLogger(__name__)

class ReportGenerator():
    def __init__(self, possible_types: List, possible_keys: List) -> None:
        self.possible_types: List = possible_types
        self.possible_keys: List = possible_keys
        self.ppt = Presentation()
        self.Title_Layout = self.ppt.slide_layouts[0]
        self.Content_Layout = self.ppt.slide_layouts[5]

    def keys_validation(self, keys: List) -> bool:
        '''
        Check whether the keys are valid.
        This is an extra validation step after the json schema validation.
        '''
        for key in keys:
            if key not in self.possible_keys:
                return False
        return True

    def generate_report(self, data: List[dict]) -> None:
        '''
        Generates slides from slide data (dictionary) if the keys and slide type (the value of "type" key) are valid.
        '''
        for raw_slide_data in data:
            are_keys_valid: bool = self.keys_validation(list(raw_slide_data.keys()))
            if are_keys_valid == False:
                logging.error("Invalid key found")
            else:
                slide_type = raw_slide_data["type"].lower()
                if slide_type not in self.possible_types:
                    logging.error("\"{}\" Slide type not supported".format(slide_type))
                else:
                    if slide_type == "title":
                        logging.info("Generating title slide....")
                        TitleSubtitleSlide.generate_title_slide(self, raw_slide_data)
                    else: 
                        if slide_type == "text":
                            logging.info("Generating text slide....")
                            TitleOnlySlide.generate_text_slide(self, raw_slide_data)
                        elif slide_type == "list":
                            logging.info("Generating list slide....")
                            TitleOnlySlide.generate_list_slide(self, raw_slide_data)
                        elif slide_type == "picture":
                            logging.info("Generating picture slide....")
                            TitleOnlySlide.generate_picture_slide(self, raw_slide_data)
                        elif slide_type == "plot":
                            logging.info("Generating plot slide....")
                            TitleOnlySlide.generate_plot_slide(self, raw_slide_data)


    @staticmethod
    def save_output(self) -> None:
        output_path = "./data/output.pptx"
        logging.info("Attempting to save result to {}".format(output_path))
        try:
            self.ppt.save(output_path)
            logging.info("Result saved to {}".format(output_path))
        except IOError:
            logging.error("Could not save result to {}".format(output_path))
            


class TitleSubtitleSlide(ReportGenerator):
    '''
    Used to generate "title" slide types.
    Uses the parent class' Title_Layout object, where the layout mode is 0.
    '''
    def __init__(self) -> None:
        super().__init__()


    def generate_title_slide(self, slide_data: dict) -> None:
        '''
        Adds title text and subtitle text on the title slide.
        '''
        self.new_slide = self.ppt.slides.add_slide(self.Title_Layout)
        self.new_slide.shapes.title.text = slide_data.get("title")
        self.new_slide.placeholders[1].text = slide_data.get("content")
        logging.info("Title Slide Done")



class TitleOnlySlide(ReportGenerator):
    '''
    Used to generate "text", "list", "picture" and "plot" slide types.
    Uses the parent class' Content_Layout object, where the layout mode is 5.
    '''
    def __init__(self) -> None:
        super().__init__()


    def generate_text_slide(self, slide_data: dict) -> None:
        '''
        Adds title text and creates textbox for text content on slide. 
        '''
        self.new_slide = self.ppt.slides.add_slide(self.Content_Layout)
        self.new_slide.shapes.title.text = slide_data.get("title")
        textbox = self.new_slide.shapes.add_textbox(Inches(1), Inches(1.5),Inches(3), Inches(1))
        textbox.text = slide_data.get("content")
        logging.info("Text Slide Done")


    def generate_list_slide(self, slide_data: dict) -> None:
        '''
        Adds title text and creates textbox and paragraphs for list content.
        '''
        self.new_slide = self.ppt.slides.add_slide(self.Content_Layout)
        self.new_slide.shapes.title.text = slide_data.get("title")

        textbox = self.new_slide.shapes.add_textbox(Inches(1), Inches(1.5),Inches(3), Inches(1))
        list_content = slide_data.get("content")
        textframe = textbox.text_frame

        for level_data in list_content:
            level_num: int = int(level_data.get("level"))
            paragraph = textframe.add_paragraph()
            if level_num == 1:
                paragraph.text = '- ' + str(level_data.get("text"))
                paragraph.font.size = Pt(30)
            elif level_num == 2:
                paragraph.text = '      ' + str(level_data.get("text"))
                paragraph.font.size = Pt(20)
            else:
                paragraph.text = '            ' + str(level_data.get("text"))
                paragraph.font.size = Pt(10)
        logging.info("List Slide Done")


    def generate_picture_slide(self, slide_data: dict) -> None:
        '''
        Adds title text and picture on slide.
        Picture must be in "data" folder
        '''
        self.new_slide = self.ppt.slides.add_slide(self.Content_Layout)
        self.new_slide.shapes.title.text = slide_data.get("title")
        left = top = Inches(1.8) 
        picture_path = "./data/" + slide_data.get("content")
        if path.isfile(picture_path):
            self.new_slide.shapes.add_picture(picture_path, left, top)
            logging.info("Picture Slide Done")
        else:
            logging.error("Picture not found in data folder - provide only the name of the file")
    

    def generate_plot_slide(self, slide_data: dict) -> None:
        '''
        Adds title text and creates XyChart. 
        Data file containing the configuration (x and y axis numbers) must be in "data" folder, separated by ';'.
        '''
        self.new_slide = self.ppt.slides.add_slide(self.Content_Layout)
        self.new_slide.shapes.title.text = slide_data.get("title")

        chart_data = XyChartData()
        series = chart_data.add_series('Series')
        series.has_title = False

        plot_file_suffix = Path(slide_data.get('content')).suffix
        plot_file_path = "./data/" + slide_data.get('content').replace(plot_file_suffix, '.csv')
        if path.isfile(plot_file_path):
            arr = np.loadtxt(plot_file_path, delimiter=';', dtype = float)
            for a in arr:
                series.add_data_point(a[0], a[1])
                # logging.info("X value: {}, Y value: {}".format(a[0], a[1]))
            
            x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
            chart = self.new_slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, x, y, cx, cy, chart_data).chart
            
            axis_labels: dict = slide_data.get("configuration")
            category_axis_title = chart.category_axis.axis_title
            category_axis_title.text_frame.text = axis_labels.get("x-label")
            value_axis_title = chart.value_axis.axis_title
            value_axis_title.text_frame.text = axis_labels.get("y-label")
            logging.info("Plot Slide Done")
        else:
            logging.error("Plot file not found in data folder - provide only the name of the file")

        

    